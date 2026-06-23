"""Export structured slide specs to a native, editable .pptx file.

Creates real DrawingML shapes (titles, text boxes, backgrounds, dividers)
so every element is clickable and editable in PowerPoint — not flat images.
"""

from __future__ import annotations

import asyncio
import json
import logging
from pathlib import Path

logger = logging.getLogger(__name__)

_CANVAS_SIZES_EMU = {
    "ppt169": (12192000, 6858000),
    "ppt43": (9144000, 6858000),
}

_CANVAS_SIZES_PX = {
    "ppt169": (1280, 720),
    "ppt43": (1024, 768),
}


def _hex_to_rgbcolor(hex_str: str):
    from pptx.dml.color import RGBColor
    h = hex_str.lstrip("#")
    if len(h) == 3:
        h = "".join(c * 2 for c in h)
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _build_pptx(design_spec: str, output_path: Path, canvas_format: str):
    from pptx import Presentation
    from pptx.util import Emu, Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    try:
        spec = json.loads(design_spec)
    except json.JSONDecodeError:
        spec = {"pages": [{"page_number": 1, "title": "Presentation", "layout": "content", "content_points": [design_spec[:500]]}]}

    theme = spec.get("theme", {})
    pages = spec.get("pages", [])
    primary_color = theme.get("primary_color", "#1e40af")
    secondary_color = theme.get("secondary_color", "#3b82f6")
    bg_color = theme.get("background_color", "#ffffff")
    accent_color = theme.get("accent_color", "#f59e0b")
    font_family = theme.get("font_family", "Microsoft YaHei")

    prs = Presentation()
    w_emu, h_emu = _CANVAS_SIZES_EMU.get(canvas_format, _CANVAS_SIZES_EMU["ppt169"])
    prs.slide_width = Emu(w_emu)
    prs.slide_height = Emu(h_emu)
    blank_layout = prs.slide_layouts[6]

    for page in pages:
        slide = prs.slides.add_slide(blank_layout)
        layout = page.get("layout", "content")
        title = page.get("title", "")
        points = page.get("content_points", [])
        notes_text = page.get("notes", "")

        # Background fill
        bg_fill = slide.background.fill
        bg_fill.solid()
        bg_fill.fore_color.rgb = _hex_to_rgbcolor(bg_color)

        if layout == "title":
            _build_title_slide(slide, title, points, primary_color, accent_color, font_family, w_emu, h_emu)
        elif layout == "closing":
            _build_closing_slide(slide, title, points, primary_color, secondary_color, font_family, w_emu, h_emu)
        elif layout == "two_column":
            _build_two_column_slide(slide, title, points, primary_color, font_family, w_emu, h_emu)
        else:
            _build_content_slide(slide, title, points, primary_color, secondary_color, accent_color, font_family, w_emu, h_emu)

        if notes_text:
            slide.notes_slide.notes_text_frame.text = notes_text

    prs.save(str(output_path))


def _add_textbox(slide, left, top, width, height, text, font_size, font_color, font_family, bold=False, alignment=None, anchor=None):
    from pptx.util import Emu, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    if anchor:
        tf.paragraphs[0].alignment = alignment

    p = tf.paragraphs[0]
    p.text = text
    run = p.runs[0] if p.runs else p.add_run()
    if not p.runs:
        run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = _hex_to_rgbcolor(font_color)
    run.font.name = font_family
    run.font.bold = bold
    if alignment:
        p.alignment = alignment
    return txBox


def _add_rect(slide, left, top, width, height, fill_color):
    from pptx.util import Emu
    from pptx.enum.shapes import MSO_SHAPE

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(left), Emu(top), Emu(width), Emu(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _hex_to_rgbcolor(fill_color)
    shape.line.fill.background()
    return shape


def _build_title_slide(slide, title, points, primary_color, accent_color, font_family, w, h):
    from pptx.enum.text import PP_ALIGN

    _add_rect(slide, 0, 0, w, h, primary_color)

    accent_bar_h = int(h * 0.008)
    _add_rect(slide, int(w * 0.15), int(h * 0.55), int(w * 0.7), accent_bar_h, accent_color)

    _add_textbox(slide, int(w * 0.1), int(h * 0.2), int(w * 0.8), int(h * 0.3),
                 title, 40, "#ffffff", font_family, bold=True, alignment=PP_ALIGN.CENTER)

    subtitle = points[0] if points else ""
    if subtitle:
        _add_textbox(slide, int(w * 0.15), int(h * 0.6), int(w * 0.7), int(h * 0.15),
                     subtitle, 20, "#e0e7ff", font_family, alignment=PP_ALIGN.CENTER)


def _build_content_slide(slide, title, points, primary_color, secondary_color, accent_color, font_family, w, h):
    from pptx.util import Emu, Pt
    from pptx.enum.text import PP_ALIGN

    _add_rect(slide, 0, 0, w, int(h * 0.005), primary_color)

    _add_textbox(slide, int(w * 0.06), int(h * 0.06), int(w * 0.88), int(h * 0.12),
                 title, 28, primary_color, font_family, bold=True)

    _add_rect(slide, int(w * 0.06), int(h * 0.17), int(w * 0.08), int(h * 0.005), accent_color)

    if points:
        body_top = int(h * 0.22)
        body_height = int(h * 0.7)
        txBox = slide.shapes.add_textbox(
            Emu(int(w * 0.06)), Emu(body_top),
            Emu(int(w * 0.88)), Emu(body_height)
        )
        tf = txBox.text_frame
        tf.word_wrap = True

        for i, pt in enumerate(points):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"  {pt}"
            p.space_after = Pt(10)
            p.space_before = Pt(4)
            for run in p.runs:
                run.font.size = Pt(16)
                run.font.color.rgb = _hex_to_rgbcolor("#374151")
                run.font.name = font_family

            bullet_shape = slide.shapes.add_shape(
                __import__("pptx.enum.shapes", fromlist=["MSO_SHAPE"]).MSO_SHAPE.OVAL,
                Emu(int(w * 0.06)),
                Emu(body_top + int(body_height * i / max(len(points), 1)) + int(h * 0.012)),
                Emu(int(w * 0.008)),
                Emu(int(w * 0.008)),
            )
            bullet_shape.fill.solid()
            bullet_shape.fill.fore_color.rgb = _hex_to_rgbcolor(secondary_color)
            bullet_shape.line.fill.background()


def _build_two_column_slide(slide, title, points, primary_color, font_family, w, h):
    from pptx.util import Emu, Pt
    from pptx.enum.text import PP_ALIGN

    _add_rect(slide, 0, 0, w, int(h * 0.005), primary_color)

    _add_textbox(slide, int(w * 0.06), int(h * 0.06), int(w * 0.88), int(h * 0.12),
                 title, 28, primary_color, font_family, bold=True)

    mid = len(points) // 2 or 1
    left_points = points[:mid]
    right_points = points[mid:]

    col_w = int(w * 0.42)
    body_top = int(h * 0.22)
    body_h = int(h * 0.7)

    for col_idx, (col_left, col_pts) in enumerate([
        (int(w * 0.06), left_points),
        (int(w * 0.52), right_points),
    ]):
        if not col_pts:
            continue
        txBox = slide.shapes.add_textbox(Emu(col_left), Emu(body_top), Emu(col_w), Emu(body_h))
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, pt in enumerate(col_pts):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"• {pt}"
            p.space_after = Pt(8)
            for run in p.runs:
                run.font.size = Pt(15)
                run.font.color.rgb = _hex_to_rgbcolor("#374151")
                run.font.name = font_family


def _build_closing_slide(slide, title, points, primary_color, secondary_color, font_family, w, h):
    from pptx.enum.text import PP_ALIGN

    _add_rect(slide, 0, 0, w, h, primary_color)

    _add_textbox(slide, int(w * 0.1), int(h * 0.3), int(w * 0.8), int(h * 0.2),
                 title or "Thank You", 44, "#ffffff", font_family, bold=True, alignment=PP_ALIGN.CENTER)

    subtitle = points[0] if points else ""
    if subtitle:
        _add_textbox(slide, int(w * 0.15), int(h * 0.55), int(w * 0.7), int(h * 0.15),
                     subtitle, 18, "#c7d2fe", font_family, alignment=PP_ALIGN.CENTER)


async def export_pptx_native(
    design_spec: str,
    workspace: Path,
    canvas_format: str = "ppt169",
) -> Path:
    """Generate a native, editable PPTX from structured design spec JSON."""
    output_path = workspace / "output.pptx"
    await asyncio.to_thread(_build_pptx, design_spec, output_path, canvas_format)
    logger.info("Native PPTX exported: %s", output_path)
    return output_path


async def export_pptx(
    svg_dir: Path,
    workspace: Path,
    canvas_format: str = "ppt169",
    design_spec: str | None = None,
) -> Path:
    """Export PPTX with priority chain:
    1. ppt-master svg_to_pptx (native DrawingML from SVG — best quality)
    2. Built-in native shapes from design_spec JSON
    3. SVG→PNG embed fallback
    """
    output_path = workspace / "output.pptx"

    # Priority 1: ppt-master's svg_to_pptx (produces real editable shapes from SVG)
    if svg_dir.exists() and any(svg_dir.glob("*.svg")):
        pptmaster_result = await _export_via_ppt_master(svg_dir, output_path, canvas_format)
        if pptmaster_result:
            return pptmaster_result

    # Priority 2: Built-in native shapes from structured JSON spec
    if design_spec:
        return await export_pptx_native(design_spec, workspace, canvas_format)

    # Priority 3: SVG→PNG image embed fallback
    return await _export_svg_fallback(svg_dir, output_path, canvas_format)


async def _export_via_ppt_master(
    svg_dir: Path, output_path: Path, canvas_format: str
) -> Path | None:
    """Use ppt-master's DrawingML converter for highest quality native PPTX."""
    import sys

    from pptx_backend.config import PPT_MASTER_SCRIPTS

    pkg_dir = PPT_MASTER_SCRIPTS / "svg_to_pptx"
    if not (pkg_dir / "pptx_builder.py").exists():
        logger.info("ppt-master svg_to_pptx package not found at %s", pkg_dir)
        return None

    scripts_str = str(PPT_MASTER_SCRIPTS)
    if scripts_str not in sys.path:
        sys.path.insert(0, scripts_str)

    try:
        from svg_to_pptx.pptx_builder import create_pptx_with_native_svg

        svg_files = sorted(svg_dir.glob("*.svg"))
        if not svg_files:
            return None

        def _run():
            return create_pptx_with_native_svg(
                svg_files=svg_files,
                output_path=output_path,
                canvas_format=canvas_format,
                use_native_shapes=True,
                use_compat_mode=True,
                verbose=False,
                transition="fade",
                transition_duration=0.4,
            )

        ok = await asyncio.to_thread(_run)
        if ok and output_path.exists():
            logger.info("ppt-master native export succeeded: %s", output_path)
            return output_path
        logger.warning("ppt-master export returned %s", ok)
    except Exception:
        logger.exception("ppt-master export error, falling back")

    return None


async def _export_svg_fallback(
    svg_dir: Path, output_path: Path, canvas_format: str
) -> Path:
    """Fallback: embed SVGs as PNG images."""
    from pptx import Presentation
    from pptx.util import Emu

    def _build():
        prs = Presentation()
        w, h = _CANVAS_SIZES_EMU.get(canvas_format, _CANVAS_SIZES_EMU["ppt169"])
        prs.slide_width = Emu(w)
        prs.slide_height = Emu(h)
        blank_layout = prs.slide_layouts[6]

        svg_files = sorted(svg_dir.glob("*.svg"))
        for svg_file in svg_files:
            slide = prs.slides.add_slide(blank_layout)
            try:
                import cairosvg
                png_path = svg_file.with_suffix(".png")
                cairosvg.svg2png(url=str(svg_file), write_to=str(png_path), output_width=1280, output_height=720)
                slide.shapes.add_picture(str(png_path), Emu(0), Emu(0), Emu(w), Emu(h))
            except ImportError:
                from pptx.util import Inches
                txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(6))
                txBox.text_frame.text = svg_file.stem

        prs.save(str(output_path))

    await asyncio.to_thread(_build)
    return output_path
